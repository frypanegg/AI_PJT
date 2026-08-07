# -*- coding: utf-8 -*-
"""발표 양식에서 필요한 슬라이드를 복제해 17장 골격을 만든다.

양식 원본 7장:
  T1 타이틀 / T2 섹션간지 / T3 목차 / T4 제목+부제 / T5 제목+본문 / T6 제목+4행바 / T7 감사합니다

목표 구성 (17장):
   1 타이틀            T1
   2 목차              T3
   3 간지 01           T2
   4 문제: 현재 상황    T6
   5 문제를 보는 관점   T5
   6 간지 02           T2
   7 시스템 개요        T5
   8 주요 기능          T6
   9 사용 기술·교육연계  T5
  10 시행착오와 학습     T6
  11 피드백 반영        T6
  12 간지 03           T2
  13 시연 영상          T5
  14 기대 효과          T5
  15 인사이트·회고      T6
  16 교육 소감          T5
  17 감사합니다         T7

add_slide.py로 복제한 뒤 <p:sldIdLst> 순서를 최종 배열로 바꾼다.
구조 작업을 먼저 끝내고 내용은 그다음에 채운다 (복제가 편집본을 복사하지 않도록).
"""

import glob
import os
import shutil
import subprocess
import sys
import zipfile

from defusedxml import minidom

SCR = os.path.dirname(os.path.abspath(__file__))
TOOLS = os.path.join(SCR, "pptx_tools", "scripts", "add_slide.py")
SRC = glob.glob(r"C:\Users\admin\Downloads\26-2*.pptx")[0]
WORK = os.path.join(SCR, "unpacked")
OUT = os.path.join(SCR, "skeleton.pptx")

# 목표 순서: 각 항목은 원본 템플릿 슬라이드 번호
TARGET = [1, 3, 2, 6, 5, 2, 5, 6, 5, 6, 6, 2, 5, 5, 6, 5, 7]

if os.path.exists(WORK):
    shutil.rmtree(WORK)
with zipfile.ZipFile(SRC) as z:
    z.extractall(WORK)
print(f"압축 해제: {len(os.listdir(os.path.join(WORK, 'ppt', 'slides')))} 항목")

# 원본 7장은 slide1..slide7. 각 템플릿이 몇 장 더 필요한지 센다.
need = {}
for t in TARGET:
    need[t] = need.get(t, 0) + 1

# 복제본 생성 — 원본 1장은 이미 있으므로 (개수-1)장씩 추가
created = {t: [f"slide{t}.xml"] for t in set(TARGET)}
for t in sorted(need):
    for _ in range(need[t] - 1):
        r = subprocess.run(
            [sys.executable, TOOLS, WORK, f"slide{t}.xml"],
            capture_output=True, text=True, encoding="utf-8",
        )
        if r.returncode != 0:
            print("복제 실패:", r.stderr[-400:]); sys.exit(1)
        # "Created ppt/slides/slideNN.xml from slideT.xml"
        newname = r.stdout.strip().split()[1].split("/")[-1]
        created[t].append(newname)
        print(f"  {t} → {newname}")

# 최종 순서대로 sldId 배열을 다시 쓴다
pres_path = os.path.join(WORK, "ppt", "presentation.xml")
rels_path = os.path.join(WORK, "ppt", "_rels", "presentation.xml.rels")

rels = minidom.parse(rels_path)
slide_rid = {}
for rel in rels.getElementsByTagName("Relationship"):
    tgt = rel.getAttribute("Target")
    if "slides/slide" in tgt:
        slide_rid[tgt.split("/")[-1]] = rel.getAttribute("Id")

doc = minidom.parse(pres_path)
lst = doc.getElementsByTagName("p:sldIdLst")[0]
for child in list(lst.childNodes):
    lst.removeChild(child)

pool = {t: list(v) for t, v in created.items()}
sid = 256
for t in TARGET:
    name = pool[t].pop(0)
    node = doc.createElement("p:sldId")
    node.setAttribute("id", str(sid))
    node.setAttribute("r:id", slide_rid[name])
    lst.appendChild(node)
    sid += 1

with open(pres_path, "w", encoding="utf-8") as f:
    f.write(doc.toxml())

# 순서에 없는 슬라이드(T4)는 clean.py가 정리한다
clean = os.path.join(SCR, "pptx_tools", "scripts", "clean.py")
r = subprocess.run([sys.executable, clean, WORK], capture_output=True, text=True, encoding="utf-8")
print(r.stdout.strip()[-300:] or r.stderr.strip()[-300:])

if os.path.exists(OUT):
    os.remove(OUT)
zf = zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED)
for cur, _, files in os.walk(WORK):
    for f in files:
        full = os.path.join(cur, f)
        zf.write(full, os.path.relpath(full, WORK).replace(os.sep, "/"))
zf.close()

from pptx import Presentation
prs = Presentation(OUT)
print(f"\n골격 완성: {OUT}  ({len(prs.slides)}장)")
for i, s in enumerate(prs.slides, 1):
    n = len(s.shapes)
    print(f"  {i:>2}. shapes={n:<3} layout={s.slide_layout.name}")
