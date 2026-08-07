# -*- coding: utf-8 -*-
"""렌더러가 없는 환경에서의 PPTX 기하 QA.

LibreOffice/PowerPoint가 모두 없어 이미지 렌더 검수를 할 수 없으므로,
스킬이 지적하는 결함 유형(경계 이탈, 여백 부족, 요소 겹침, 텍스트 넘침)을
좌표와 폰트 메트릭 추정으로 대신 점검한다.
"""

import sys

from pptx import Presentation
from pptx.util import Emu

PATH = sys.argv[1]
EMU_IN = 914400.0

prs = Presentation(PATH)

# 슬라이드 크기는 문서마다 다르다 (16:9 13.3" 도 있고 26.67" 대형 캔버스도 있다).
# 고정값을 쓰면 대형 캔버스에서 전부 오탐이 나므로 실제 크기를 읽고,
# 여백 기준도 폭에 비례해 잡는다 (13.3"에서 0.5" 에 해당하는 비율).
SLIDE_W = prs.slide_width / EMU_IN
SLIDE_H = prs.slide_height / EMU_IN
MARGIN_MIN = SLIDE_W * 0.0375
print(f"슬라이드 크기: {SLIDE_W:.2f} x {SLIDE_H:.2f} in  (여백 기준 {MARGIN_MIN:.2f}\")")


def inches(v):
    return (v or 0) / EMU_IN


def glyph_width_pt(ch, size):
    """Malgun Gothic 기준 대략적 글리프 폭(pt)."""
    o = ord(ch)
    if 0xAC00 <= o <= 0xD7A3 or 0x3130 <= o <= 0x318F or 0x4E00 <= o <= 0x9FFF:
        return size * 1.0          # 한글/한자: 전각
    if ch == " ":
        return size * 0.26
    if ch in "·—–…":
        return size * 0.9
    if ch.isdigit() or ch.isupper():
        return size * 0.58
    return size * 0.50             # 소문자/기호


def measure(text, size, bold=False):
    w = sum(glyph_width_pt(c, size) for c in text)
    return w * (1.04 if bold else 1.0)


issues = {"bounds": [], "margin": [], "overflow": [], "overlap": []}

for idx, slide in enumerate(prs.slides, start=1):
    text_boxes = []

    for sh in slide.shapes:
        x, y = inches(sh.left), inches(sh.top)
        w, h = inches(sh.width), inches(sh.height)
        r, b = x + w, y + h
        name = sh.shape_type

        # 1) 슬라이드 경계 이탈
        if x < -0.02 or y < -0.02 or r > SLIDE_W + 0.02 or b > SLIDE_H + 0.02:
            # 의도적 장식(원/사각 배경)은 배경 장식으로 분류
            issues["bounds"].append(
                f"슬라이드 {idx}: {name} 위치 ({x:.2f},{y:.2f})-({r:.2f},{b:.2f})"
            )

        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue

        # 폰트 크기 수집
        sizes, bolds = [], []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(run.font.size.pt)
                bolds.append(bool(run.font.bold))
        size = max(sizes) if sizes else 12.0
        bold = any(bolds)

        text_boxes.append((x, y, r, b, txt[:34], size))

        # 2) 여백 (장식/배경 텍스트 제외 — 본문 텍스트만)
        if x < MARGIN_MIN - 0.06 or r > SLIDE_W - MARGIN_MIN + 0.06:
            issues["margin"].append(
                f"슬라이드 {idx}: '{txt[:26]}' x={x:.2f} r={r:.2f}"
            )

        # 3) 텍스트 넘침 추정
        inner_w_pt = max(w * 72 - 4, 10)   # margin:0 기준 약간의 여유
        line_h = size * 1.42
        lines = 0
        for para in sh.text_frame.paragraphs:
            ptext = "".join(run.text for run in para.runs)
            if not ptext:
                lines += 1
                continue
            need = measure(ptext, size, bold)
            lines += max(1, int(need / inner_w_pt) + (1 if need % inner_w_pt else 0))
        need_h = lines * line_h / 72.0
        if need_h > h + 0.06:
            issues["overflow"].append(
                f"슬라이드 {idx}: '{txt[:30]}' 필요 {need_h:.2f}\" > 박스 {h:.2f}\" ({size}pt, {lines}줄)"
            )

    # 4) 텍스트끼리 겹침
    for i in range(len(text_boxes)):
        for j in range(i + 1, len(text_boxes)):
            a, bb = text_boxes[i], text_boxes[j]
            ox = min(a[2], bb[2]) - max(a[0], bb[0])
            oy = min(a[3], bb[3]) - max(a[1], bb[1])
            if ox > 0.06 and oy > 0.06:
                issues["overlap"].append(
                    f"슬라이드 {idx}: '{a[4]}' ↔ '{bb[4]}' 겹침 {ox:.2f}x{oy:.2f}\""
                )

print(f"슬라이드 수: {len(prs.slides)}")
for k, label in [
    ("bounds", "슬라이드 경계 이탈"),
    ("margin", "여백 0.5\" 미만"),
    ("overflow", "텍스트 넘침(추정)"),
    ("overlap", "텍스트 겹침"),
]:
    v = issues[k]
    print(f"\n=== {label}: {len(v)}건 ===")
    for line in v[:25]:
        print("  -", line)
    if len(v) > 25:
        print(f"  ... 외 {len(v)-25}건")
